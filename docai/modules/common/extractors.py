import csv
import io
import os

from .models import AttachedFile

from ...logging_config import get_logger, log_call

logger = get_logger(__name__)

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    logger.debug("python-docx not installed — extract_word() unavailable.")
    HAS_DOCX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    logger.debug("openpyxl not installed — extract_excel() unavailable.")
    HAS_XLSX = False

try:
    import fitz
    HAS_PDF_TEXT = True
except ImportError:
    logger.debug("pymupdf (fitz) not installed — PDF text extraction unavailable.")
    HAS_PDF_TEXT = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    logger.debug("python-pptx not installed — extract_pptx() unavailable.")
    HAS_PPTX = False


@log_call
def extract_word(path: str) -> AttachedFile:
    doc = DocxDocument(path)
    lines = []
    body = doc.element.body
    for child in body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "p":
            from docx.text.paragraph import Paragraph
            paragraph = Paragraph(child, doc)
            lines.append(paragraph.text)
        elif tag == "tbl":
            from docx.table import Table
            table = Table(child, doc)
            for row in table.rows:
                cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                lines.append(" | ".join(cells))
            lines.append("")
    return AttachedFile(path=path, name=os.path.basename(path),
                        file_type="word", claude_content="\n".join(lines))


@log_call
def extract_excel(path: str) -> AttachedFile:
    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for sname in wb.sheetnames:
        ws = wb[sname]
        rows = [[str(cell) if cell is not None else "" for cell in row]
                for row in ws.iter_rows(values_only=True)]
        while rows and all(cell == "" for cell in rows[-1]):
            rows.pop()
        parts.append(f"--- Sheet: {sname} ---")
        buf = io.StringIO()
        csv.writer(buf).writerows(rows)
        parts.extend([buf.getvalue().rstrip(), ""])
    return AttachedFile(path=path, name=os.path.basename(path),
                        file_type="excel", claude_content="\n".join(parts))


@log_call
def extract_pptx(path: str) -> AttachedFile:
    prs = Presentation(path)
    parts = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {slide_index} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text.strip():
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                    parts.append(" | ".join(cells))
        parts.append("")
    return AttachedFile(path=path, name=os.path.basename(path),
                        file_type="ppt", claude_content="\n".join(parts))


@log_call
def pdf_page_texts(path: str) -> list[str]:
    """Text for each PDF page (empty if the page has no text layer — e.g. a
    pure scanned-image PDF). Used by the "Extract full text" feature (local,
    no AI) and PDF keyword search."""
    doc = fitz.open(path)
    try:
        return [page.get_text().strip() for page in doc]
    finally:
        doc.close()


@log_call
def extract_pdf(path: str) -> AttachedFile:
    pages = pdf_page_texts(path)
    parts = []
    for page_number, text in enumerate(pages, start=1):
        parts.append(f"--- Trang {page_number} ---")
        parts.append(text or "(trang trống hoặc không có lớp chữ — có thể là PDF quét ảnh)")
    return AttachedFile(path=path, name=os.path.basename(path),
                        file_type="pdf", claude_content="\n".join(parts))
