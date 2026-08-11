"""Build a new Word/Excel/PowerPoint file from AI-generated content (the
"Create a new document via chat" feature).

The AI (server) is asked to return text with `### Trang N: <title>` markers
so it can be split into pages/slides — if the AI doesn't follow the format,
`parse_sections()` falls back to splitting paragraphs evenly. Excel reuses
the existing `--- Sheet: X ---` convention already in `savers.save_excel`, no
separate syntax needed.
"""
import re
from dataclasses import dataclass

from ...logging_config import get_logger, log_call
from ...strings import Creators as S

logger = get_logger(__name__)

try:
    from docx import Document as DocxDocument
except ImportError:
    logger.debug("python-docx not installed — create_word() will raise if called.")
    DocxDocument = None  # type: ignore

try:
    import pptx
    from pptx import Presentation
except ImportError:
    logger.debug("python-pptx not installed — create_pptx() will raise if called.")
    pptx = None  # type: ignore
    Presentation = None  # type: ignore

from .savers import save_excel

HAS_PPTX = pptx is not None

_MARKER_RE = re.compile(r'^#{1,3}\s*Trang\s+(\d+)\s*:\s*(.+?)\s*$', re.IGNORECASE | re.MULTILINE)


@dataclass
class SectionSpec:
    heading: str
    body: str


@log_call
def parse_sections(raw_text: str, expected_count: int = 0) -> list[SectionSpec]:
    """Split the AI-returned text into pages/sections by the `### Trang N: …`
    marker.

    If the AI doesn't follow the format, split paragraphs evenly into about
    `expected_count` sections (just to get a reasonable result, doesn't need
    to be exactly the right page count).
    """
    raw_text = (raw_text or "").strip()
    if not raw_text:
        return []

    matches = list(_MARKER_RE.finditer(raw_text))
    if matches:
        sections = []
        for match_index, match in enumerate(matches):
            start = match.end()
            end = matches[match_index + 1].start() if match_index + 1 < len(matches) else len(raw_text)
            body = raw_text[start:end].strip()
            sections.append(SectionSpec(heading=match.group(2).strip(), body=body))
        return sections

    paragraphs = [paragraph.strip() for paragraph in re.split(r'\n\s*\n', raw_text) if paragraph.strip()]
    if not paragraphs:
        paragraphs = [raw_text]
    section_count = max(1, expected_count or 1)
    chunk = max(1, -(-len(paragraphs) // section_count))  # ceil(len/n)
    sections = []
    for start_index in range(0, len(paragraphs), chunk):
        group = paragraphs[start_index:start_index + chunk]
        sections.append(SectionSpec(heading=f"Trang {len(sections) + 1}", body="\n\n".join(group)))
    return sections


@log_call
def create_word(sections: list[SectionSpec], out_path: str):
    if DocxDocument is None:
        raise RuntimeError(S.MISSING_DOCX_LIB)
    doc = DocxDocument()
    for section_index, sec in enumerate(sections):
        if section_index > 0:
            doc.add_page_break()
        doc.add_heading(sec.heading, level=1)
        for line in sec.body.split("\n"):
            if line.strip():
                doc.add_paragraph(line.strip())
    doc.save(out_path)


@log_call
def create_pptx(sections: list[SectionSpec], out_path: str):
    if Presentation is None:
        raise RuntimeError(S.MISSING_PPTX_LIB)
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    for sec in sections:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = sec.heading
        lines = [line.strip() for line in sec.body.split("\n") if line.strip()] or [""]
        body = slide.placeholders[1].text_frame
        body.text = lines[0]
        for line in lines[1:]:
            paragraph = body.add_paragraph()
            paragraph.text = line
    prs.save(out_path)


@log_call
def create_excel_from_text(raw_text: str, out_path: str):
    save_excel(raw_text, out_path)
